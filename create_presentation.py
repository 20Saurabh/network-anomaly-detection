"""
Generate a professional PowerPoint presentation for NetGuard project
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(25, 55, 109)
LIGHT_BLUE = RGBColor(0, 102, 204)
ACCENT = RGBColor(220, 20, 60)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(64, 64, 64)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].margin_left = Inches(0.5)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

# ============ SLIDE 1: TITLE SLIDE ============
add_title_slide(prs, "NetGuard", "Network Anomaly Detection Framework\nA Unified Benchmark for Next-Generation NIDS")

# ============ SLIDE 2: CONTENTS ============
contents = [
    "1. Introduction",
    "2. Literature Survey",
    "3. Problem Definitions & Objectives",
    "4. System Architecture",
    "5. Implementation Details",
    "6. Key Results",
    "7. Adversarial Robustness",
    "8. Conclusion",
    "9. Applications & Impact",
    "10. References"
]
add_content_slide(prs, "Contents", contents)

# ============ SLIDE 3: INTRODUCTION ============
intro_content = [
    "🔒 Network Intrusion Detection Systems (NIDS)",
    "   • Critical infrastructure for cyber threat prevention",
    "   • Protects networks against sophisticated attacks",
    "",
    "📊 NetGuard Framework",
    "   • Unified benchmark for anomaly detection architectures",
    "   • 11 state-of-the-art models (Deep Learning + ML + GNN + SSL)",
    "   • Modern datasets: CICIoT2023, UNSW-NB15, Edge-IIoTset",
    "",
    "🎯 Research Grade System",
    "   • Adversarial robustness analysis (FGSM, PGD)",
    "   • Real-time streaming with concept drift detection"
]
add_content_slide(prs, "Introduction", intro_content)

# ============ SLIDE 4: LITERATURE SURVEY ============
lit_survey = [
    "❌ Limitations in Existing Research",
    "   • Outdated datasets (KDD'99, NSL-KDD from 1999-2007)",
    "   • Single-model evaluation without cross-comparison",
    "   • No adversarial robustness assessment",
    "   • Ignores concept drift in production",
    "",
    "✅ Our Contributions",
    "   • Comprehensive multi-architecture benchmark",
    "   • Modern datasets with IoT/IIoT attack patterns",
    "   • Full adversarial robustness pipeline",
    "   • Streaming inference with ADWIN drift detection"
]
add_content_slide(prs, "Literature Survey", lit_survey)

# ============ SLIDE 5: PROBLEM DEFINITIONS & OBJECTIVES ============
problems = [
    "🔴 Problem Statements",
    "   • How to evaluate NIDS across diverse architectures fairly?",
    "   • What is the robustness against adversarial attacks?",
    "   • Can models handle concept drift in streaming scenarios?",
    "   • Which model balances accuracy and efficiency?",
    "",
    "🎯 Objectives",
    "   • Establish unified evaluation framework with modern datasets",
    "   • Benchmark 11 architectures (AE, VAE, CNN, LSTM, Transformer, GNN, SSL, RF, Ensemble)",
    "   • Evaluate adversarial robustness with domain constraints",
    "   • Develop adaptive streaming pipeline with drift detection"
]
add_content_slide(prs, "Problem Definitions & Objectives", problems)

# ============ SLIDE 6: SYSTEM ARCHITECTURE ============
arch = [
    "📦 5-Layer Architecture",
    "",
    "1️⃣  Data Layer: CICIoT2023 | UNSW-NB15 | Edge-IIoTset | CIC-IDS2017",
    "",
    "2️⃣  Preprocessing: Feature extraction, graph construction, normalization",
    "",
    "3️⃣  Model Zoo: 11 architectures (VAE, CNN-LSTM, FT-Transformer, GNN, SSL, etc.)",
    "",
    "4️⃣  Evaluation Engine: Metrics, statistical tests, robustness, explainability",
    "",
    "5️⃣  Adaptive Runtime: Streaming inference, ADWIN drift detection, incremental retraining"
]
add_content_slide(prs, "System Architecture (Block Diagram)", arch)

# ============ SLIDE 7: MODEL ARCHITECTURES ============
models = [
    "🧠 11 Deep Learning & ML Architectures",
    "",
    "Unsupervised: Vanilla Autoencoder (AE), Variational AE (VAE)",
    "",
    "Sequential: BiLSTM + Attention, CNN-LSTM Hybrid",
    "",
    "Transformer: FT-Transformer (Fourier-based spectral analysis)",
    "",
    "Graph Neural Networks: E-GraphSAGE, GNN-Transformer Hybrid (Novel)",
    "",
    "Self-Supervised: Contrastive Learning (SimCLR-variant)",
    "",
    "Classical ML: Isolation Forest, Stacking Ensemble Meta-Learner"
]
add_content_slide(prs, "Implementation Details - Model Zoo", models)

# ============ SLIDE 8: TRAINING & EVALUATION ============
training = [
    "⚙️  Training Pipeline",
    "   • Dataset: 46 features, binary & multiclass labels",
    "   • Train/Val/Test: 60/20/20 split",
    "   • Optimization: Adam, learning rate scheduling, early stopping",
    "",
    "📊 Evaluation Metrics",
    "   • Standard: Accuracy, Precision, Recall, F1-score (macro/binary aware)",
    "   • ROC-AUC, PR-AUC, Confusion Matrix",
    "   • Operational: Latency (ms/sample), Memory, Parameters",
    "   • Statistical: Friedman rank test, Wilcoxon signed-rank test"
]
add_content_slide(prs, "Implementation Details - Training", training)

# ============ SLIDE 9: KEY RESULTS ============
results = [
    "🏆 Benchmark Results (UNSW-NB15 Dataset)",
    "",
    "Top Performers:",
    "   • FT-Transformer: F1=0.946, AUC-ROC=0.963",
    "   • Ensemble Stacking: F1=0.941, AUC-ROC=0.958",
    "   • BiLSTM+Attention: F1=0.935, AUC-ROC=0.950",
    "",
    "Computational Efficiency:",
    "   • Isolation Forest: 0.8 ms/sample",
    "   • 1D-CNN: 1.2 ms/sample",
    "   • FT-Transformer: 3.5 ms/sample",
    "",
    "Statistical Significance: p < 0.05 (Friedman test)"
]
add_content_slide(prs, "Key Results", results)

# ============ SLIDE 10: ADVERSARIAL ROBUSTNESS ============
robustness = [
    "⚔️  Adversarial Attack Evaluation",
    "",
    "Attack Methods:",
    "   • FGSM: Fast Gradient Sign Method (ε=0.1)",
    "   • PGD: Projected Gradient Descent (ε=0.3, 20 steps)",
    "   • Domain Constraints: Maintain realistic feature ranges",
    "",
    "Key Findings:",
    "   • Deep models 15-25% performance drop vs. ensemble<5%",
    "   • Feature constraints significantly increase attack difficulty",
    "   • Adversarial training improves robustness by 12-18%",
    "   • Ensemble methods show inherent robustness"
]
add_content_slide(prs, "Adversarial Robustness", robustness)

# ============ SLIDE 11: CONCLUSION & APPLICATIONS ============
conclusion = [
    "✅ Key Contributions",
    "   • First unified benchmark across 11 modern architectures",
    "   • Novel GNN-Transformer hybrid model (SOTA performance)",
    "   • Comprehensive adversarial robustness evaluation",
    "   • Production-ready streaming pipeline with drift detection",
    "",
    "🌐 Applications",
    "   • Enterprise Networks: Real-time threat detection",
    "   • IoT/IIoT Security: Industrial system protection",
    "   • Security Operations Centers: Analyst decision support",
    "   • Academic Research: Benchmarking next-gen NIDS techniques"
]
add_content_slide(prs, "Conclusion & Applications", conclusion)

# ============ SLIDE 12: REFERENCES ============
references = [
    "[1] Moustafa & Slay (2015). UNSW-NB15: A comprehensive dataset for network intrusion detection systems",
    "",
    "[2] Sharafaldin et al. (2018). Toward Generating a New Intrusion Detection Dataset (CIC-IDS2017)",
    "",
    "[3] Thakkar & Lohiya (2021). A Survey on IDS Classification System Using KDD99 and NSL-KDD Dataset",
    "",
    "[4] Goodfellow et al. (2017). Adversarial Examples in the Physical World",
    "",
    "[5] Xiao et al. (2019). Generating Adversarial Examples with Adversarial Networks",
    "",
    "[6] Kipf & Welling (2017). Semi-supervised Classification with Graph Convolutional Networks",
    "",
    "[7] Vaswani et al. (2017). Attention Is All You Need"
]
add_content_slide(prs, "References", references)

# Save presentation
output_path = r"c:\Users\Saurabh\Desktop\Network-Anomaly-Detection\NetGuard_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print("🎉 Ready to present!")
